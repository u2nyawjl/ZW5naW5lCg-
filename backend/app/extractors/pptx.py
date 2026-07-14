import io

from pptx import Presentation


def extract(content: bytes) -> tuple[str, list[str]]:
    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        lines = [f"## Diapositiva {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            lines.append(f"> Notas: {slide.notes_slide.notes_text_frame.text.strip()}")
        if len(lines) > 1:
            parts.append("\n".join(lines))

    return "\n\n".join(parts), []


def metadata(content: bytes) -> dict:
    props = Presentation(io.BytesIO(content)).core_properties
    return {
        "author": props.author,
        "title": props.title,
        "created": props.created.isoformat() if props.created else None,
        "modified": props.modified.isoformat() if props.modified else None,
        "last_modified_by": props.last_modified_by,
        "revision": props.revision,
    }
