"""Fixtures generados en runtime.

El string EICAR se ensambla por partes a propósito: escrito literal en el repo, el
antivirus del host pondría en cuarentena el propio archivo de test.
"""

import io
import zipfile


def eicar() -> bytes:
    parts = ["X5O!P%@AP[4\\PZX54(P^)7CC)7}$EICAR-STANDARD-", "ANTIVIRUS-TEST-FILE!$H+H*"]
    return "".join(parts).encode()


def minimal_pdf(text: str = "Informe de capstone U2NyaWJl") -> bytes:
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + body + b"\nendobj\n"

    xref_at = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_at}\n%%EOF"
    ).encode()
    return bytes(out)


def minimal_xlsx() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Telemetria"
    ws.append(["mac_address", "rssi", "visto_en"])
    ws.append(["AA:BB:CC:DD:EE:FF", -67, "2026-07-14"])
    wb.properties.creator = "Nico"

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def minimal_pptx() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Avance capstone"
    slide.placeholders[1].text = "Pipeline de seguridad operativo"
    prs.core_properties.author = "Nico"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def windows_executable() -> bytes:
    """Cabecera PE: lo que de verdad hay dentro de un 'factura.pdf' malicioso."""
    return b"MZ\x90\x00\x03\x00\x00\x00\x04\x00\x00\x00\xff\xff\x00\x00" + b"\x00" * 128


def zip_bomb(uncompressed_mb: int = 8) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("xl/worksheets/sheet1.xml", b"\x00" * (uncompressed_mb * 1_048_576))
    return buf.getvalue()
